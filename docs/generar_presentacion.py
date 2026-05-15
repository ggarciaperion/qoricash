"""
QoriCash FX — Presentación Corporativa (versión minimalista)
Diseño: fondo blanco / gris muy claro, verde institucional como acento.
Navy solo en portada, cierre y elementos de contraste puntual.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
from lxml import etree
import math, random

# ── Paleta ────────────────────────────────────────────────────────────────────
GREEN      = RGBColor(0x22, 0xC5, 0x5E)   # #22C55E  primary
GREEN_D    = RGBColor(0x16, 0xA3, 0x4A)   # #16A34A  primary-600
GREEN_XL   = RGBColor(0xF0, 0xFD, 0xF4)   # #f0fdf4  primary-50
GREEN_100  = RGBColor(0xDC, 0xFC, 0xE7)   # #dcfce7  primary-100
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # #0D1B2A  secondary
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE  = RGBColor(0xF9, 0xFA, 0xFB)   # fondo slides
CARD_BG    = RGBColor(0xF3, 0xF4, 0xF6)   # #F3F4F6  gris tarjeta
BORDER     = RGBColor(0xE5, 0xE7, 0xEB)   # #E5E7EB  borde sutil
BODY       = RGBColor(0x1F, 0x29, 0x37)   # texto principal
MUTED      = RGBColor(0x6B, 0x72, 0x80)   # texto secundario

LOGO     = "/Users/gianpierre/Desktop/Qoricash/Sistema/qoricashweb/public/logo-principal.png"
OUT      = "/Users/gianpierre/Desktop/Qoricash/Sistema/qoricash/docs/QoriCash_Presentacion_Corporativa.pptx"
IMGS     = "/Users/gianpierre/Desktop/Qoricash/Sistema/qoricash/docs/imgs/"
I_FOREX    = IMGS + "forex.jpg"
I_BUILDING = IMGS + "building.jpg"
I_MEETING  = IMGS + "meeting.jpg"
I_EXCHANGE = IMGS + "exchange.jpg"
I_DIGITAL  = IMGS + "digital.jpg"
I_SECURITY = IMGS + "security.jpg"

W = Inches(13.33)
H = Inches(7.50)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ── Primitivas ────────────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill=None, line=None, lw=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = lw
    else:
        s.line.fill.background()
    return s


def tb(slide, text, x, y, w, h, size=12, bold=False,
       color=BODY, align=PP_ALIGN.LEFT, wrap=True, font="Calibri"):
    t = slide.shapes.add_textbox(x, y, w, h)
    tf = t.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return t


def logo(slide, x, y, h):
    try: return slide.shapes.add_picture(LOGO, x, y, height=h)
    except: return None


def img(slide, path, x, y, w, h, alpha_pct=100):
    """Inserta imagen. Si alpha_pct < 100 aplica transparencia vía XML."""
    try:
        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
        if alpha_pct < 100:
            # Aplica transparencia al blipFill
            blip = pic._element.find('.//' + qn('a:blip'))
            if blip is not None:
                alphaModFix = etree.SubElement(blip, qn('a:alphaModFix'))
                alphaModFix.set('amt', str(int(alpha_pct * 1000)))
        return pic
    except Exception as e:
        return None


def img_panel(slide, path, x, y, w, h, overlay_alpha=45):
    """
    Imagen recortada en un panel con overlay navy semitransparente encima,
    para que el texto sobre ella sea legible.
    """
    try:
        pic = slide.shapes.add_picture(path, x, y, width=w, height=h)
    except:
        return
    # Overlay oscuro encima para legibilidad
    ov = rect(slide, x, y, w, h, fill=NAVY)
    set_shape_fill_alpha(ov, overlay_alpha)


def icon_circle(slide, symbol, x, y, size, bg=GREEN, fg=WHITE, font_size=None):
    """Círculo de color con símbolo/letra centrado — icono minimalista."""
    circle = slide.shapes.add_shape(9, x, y, size, size)  # 9 = oval
    circle.fill.solid(); circle.fill.fore_color.rgb = bg
    circle.line.fill.background()
    fs = font_size or int(Pt(size.inches * 18))
    tb(slide, symbol,
       x, y + size * 0.12, size, size * 0.75,
       size=max(10, int(size.inches * 16)),
       bold=True, color=fg, align=PP_ALIGN.CENTER)


def icon_box(slide, symbol, x, y, size, bg=GREEN_XL, fg=GREEN_D, font_size=14, font="Segoe UI Emoji"):
    """Cuadrado redondeado con símbolo/emoji — icono tarjeta."""
    box = slide.shapes.add_shape(5, x, y, size, size)  # 5 = rounded rect
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.color.rgb = GREEN; box.line.width = Pt(0.5)
    try: box.adjustments[0] = 0.15
    except: pass
    tb(slide, symbol, x, y + size * 0.10, size, size * 0.80,
       size=font_size, bold=False, color=fg, align=PP_ALIGN.CENTER, font=font)


# ── Transparencia vía XML (alpha 0-100) ──────────────────────────────────────

def _set_alpha(color_el, alpha_pct):
    """Inserta a:alpha en un elemento a:sRgbClr o a:prstClr"""
    val = str(int(alpha_pct * 1000))
    for old in color_el.findall(qn('a:alpha')):
        color_el.remove(old)
    el = etree.SubElement(color_el, qn('a:alpha'))
    el.set('val', val)


def set_shape_fill_alpha(shape, alpha_pct):
    srgb = shape._element.find('.//' + qn('a:solidFill') + '/' + qn('a:sRgbClr'))
    if srgb is not None:
        _set_alpha(srgb, alpha_pct)


def set_line_alpha(connector, alpha_pct):
    ln = connector._element.find('.//' + qn('a:ln'))
    if ln is None: return
    srgb = ln.find('.//' + qn('a:sRgbClr'))
    if srgb is not None:
        _set_alpha(srgb, alpha_pct)


def fline(slide, x1, y1, x2, y2, color, width_pt=0.75, alpha=100):
    """Dibuja una línea con transparencia opcional"""
    cn = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, x1, y1, x2, y2)
    cn.line.color.rgb = color
    cn.line.width = Pt(width_pt)
    if alpha < 100:
        set_line_alpha(cn, alpha)
    return cn


# ── Fondo forex para la portada ──────────────────────────────────────────────

def forex_cover_bg(slide):
    """
    Fondo forex minimalista: solo grid tenue + una curva de precio suave.
    Todo al mínimo de opacidad para no competir con el texto.
    """
    random.seed(7)

    # Grid horizontal — 7 líneas, 3% opacidad
    for i in range(1, 8):
        gy = H * i / 8
        fline(slide, Emu(0), gy, W, gy, GREEN, width_pt=0.35, alpha=3)

    # Grid vertical — 9 líneas, 2% opacidad
    for i in range(1, 10):
        gx = W * i / 10
        fline(slide, gx, Emu(0), gx, H, GREEN, width_pt=0.35, alpha=2)

    # Curva de precio única, suave, 7% opacidad
    n = 60
    pts = []
    y_frac = 0.55
    for i in range(n):
        drift = (0.52 - y_frac) * 0.06
        y_frac += random.uniform(-0.018, 0.016) + drift
        y_frac = max(0.28, min(0.78, y_frac))
        pts.append((i / (n - 1), y_frac))

    for i in range(len(pts) - 1):
        x1 = Inches(0.50) + int((W - Inches(0.50)) * pts[i][0])
        y1 = int(H * pts[i][1])
        x2 = Inches(0.50) + int((W - Inches(0.50)) * pts[i + 1][0])
        y2 = int(H * pts[i + 1][1])
        fline(slide, x1, y1, x2, y2, GREEN, width_pt=0.8, alpha=7)


def content_bg(slide):
    """
    Fondo blanco con textura forex sutil:
      — Grid horizontal + vertical muy tenue
      — Curva de precio suave
      — Símbolos de moneda watermark ($, S/, USD, PEN)
      — Micro-velas (candlesticks) decorativos
    Todo a opacidad mínima para no competir con el contenido.
    """
    rect(slide, 0, 0, W, H, fill=WHITE)

    # ── Textura forex ─────────────────────────────────────────────────────
    random.seed(21)
    GHOST = RGBColor(0xC6, 0xF0, 0xD4)   # verde muy pálido para texto watermark

    # Grid horizontal — 6 líneas, alpha 4 %
    for i in range(1, 7):
        gy = H * i / 7
        fline(slide, Emu(0), gy, W, gy, GREEN, width_pt=0.25, alpha=4)

    # Grid vertical — 9 líneas, alpha 3 %
    for i in range(1, 10):
        gx = W * i / 10
        fline(slide, gx, Emu(0), gx, H, GREEN, width_pt=0.25, alpha=3)

    # Curva de precio principal — alpha 6 %
    n = 60
    pts = []
    yf = 0.56
    for i in range(n):
        drift = (0.52 - yf) * 0.05
        yf += random.uniform(-0.014, 0.013) + drift
        yf = max(0.22, min(0.84, yf))
        pts.append((i / (n - 1), yf))
    for i in range(len(pts) - 1):
        x1 = int(W * pts[i][0]);     y1 = int(H * pts[i][1])
        x2 = int(W * pts[i + 1][0]); y2 = int(H * pts[i + 1][1])
        fline(slide, x1, y1, x2, y2, GREEN, width_pt=0.55, alpha=6)

    # Segunda curva secundaria — alpha 3 %
    random.seed(55)
    pts2 = []
    yf2 = 0.38
    for i in range(n):
        drift2 = (0.42 - yf2) * 0.04
        yf2 += random.uniform(-0.011, 0.011) + drift2
        yf2 = max(0.15, min(0.65, yf2))
        pts2.append((i / (n - 1), yf2))
    for i in range(len(pts2) - 1):
        x1 = int(W * pts2[i][0]);     y1 = int(H * pts2[i][1])
        x2 = int(W * pts2[i + 1][0]); y2 = int(H * pts2[i + 1][1])
        fline(slide, x1, y1, x2, y2, GREEN, width_pt=0.35, alpha=3)

    # Micro-velas decorativas — alpha 5 %
    random.seed(33)
    candle_positions = [
        (0.08, 0.30), (0.18, 0.55), (0.28, 0.42), (0.38, 0.65),
        (0.62, 0.38), (0.72, 0.58), (0.82, 0.45), (0.92, 0.62),
    ]
    for (cxf, cyf) in candle_positions:
        ch = random.uniform(0.05, 0.10)
        cw = Inches(0.10)
        cx = int(W * cxf); cy = int(H * cyf); cH = int(H * ch)
        body = rect(slide, cx, cy, cw, cH, fill=GREEN)
        set_shape_fill_alpha(body, 5)
        body.line.fill.background()
        # Mecha superior
        fline(slide, cx + int(cw / 2), cy - int(H * 0.018),
              cx + int(cw / 2), cy, GREEN, width_pt=0.4, alpha=4)
        # Mecha inferior
        fline(slide, cx + int(cw / 2), cy + cH,
              cx + int(cw / 2), cy + cH + int(H * 0.018), GREEN, width_pt=0.4, alpha=4)

    # Símbolos de moneda watermark
    sym_data = [
        ("$",    0.03,  0.10, 52),
        ("S/",   0.84,  0.06, 44),
        ("USD",  0.88,  0.74, 28),
        ("PEN",  0.04,  0.78, 28),
        ("$",    0.46,  0.84, 38),
        ("S/",   0.36,  0.12, 30),
        ("1.0",  0.66,  0.88, 22),
        ("3.5",  0.20,  0.42, 22),
        ("USD",  0.54,  0.20, 20),
        ("PEN",  0.74,  0.32, 20),
    ]
    for sym, xf, yf_s, fs in sym_data:
        t = slide.shapes.add_textbox(
            int(W * xf), int(H * yf_s), Inches(1.40), Inches(0.65))
        tf = t.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = sym
        r.font.size = Pt(fs)
        r.font.bold = True
        r.font.color.rgb = GHOST
        r.font.name = "Calibri"

    # ── Elementos fijos encima ─────────────────────────────────────────────
    rect(slide, 0, 0, W, Inches(0.07), fill=GREEN)
    rect(slide, 0, H - Inches(0.06), W, Inches(0.06), fill=BORDER)
    logo(slide, W - Inches(1.10), H - Inches(0.58), Inches(0.42))


def section_title(slide, title, subtitle=None):
    """Encabezado: acento verde izquierdo + título navy"""
    content_bg(slide)
    rect(slide, Inches(0.40), Inches(0.22), Inches(0.06), Inches(0.70), fill=GREEN)
    tb(slide, title,
       Inches(0.60), Inches(0.20), W - Inches(1.0), Inches(0.60),
       size=26, bold=True, color=NAVY)
    if subtitle:
        tb(slide, subtitle,
           Inches(0.60), Inches(0.80), W - Inches(1.0), Inches(0.35),
           size=12, color=MUTED)
    rect(slide, Inches(0.60), Inches(1.15), Inches(12.10), Inches(0.015), fill=BORDER)
    return Inches(1.30)   # y_content start


def card(slide, x, y, w, h, accent=True):
    """Tarjeta blanca con sombra simulada (borde gris) y acento verde top opcional"""
    rect(slide, x, y, w, h, fill=WHITE, line=BORDER, lw=Pt(0.5))
    if accent:
        rect(slide, x, y, w, Inches(0.08), fill=GREEN)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — PORTADA  (navy — la única slide oscura)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=NAVY)
# Imagen edificio corporativo en el lado derecho — muy tenue
img(sl, I_BUILDING, Inches(6.80), Inches(0), Inches(6.53), H, alpha_pct=12)
# Fondo tecnológico forex (grid + curvas)
forex_cover_bg(sl)
# Franja verde izquierda estrecha
rect(sl, 0, 0, Inches(0.50), H, fill=GREEN)
# Acento verde inferior
rect(sl, Inches(0.50), H - Inches(0.10), W - Inches(0.50), Inches(0.10), fill=GREEN_D)

logo(sl, Inches(1.10), Inches(1.00), Inches(1.90))

tb(sl, "QoriCash FX",
   Inches(3.30), Inches(1.10), Inches(8.50), Inches(1.10),
   size=50, bold=True, color=WHITE)
tb(sl, "QORICASH S.A.C.  ·  RUC: 20615113698",
   Inches(3.30), Inches(2.25), Inches(8.50), Inches(0.40),
   size=13, color=GREEN)
tb(sl, "Tu casa de cambio digital de confianza",
   Inches(1.10), Inches(3.40), Inches(11.0), Inches(0.65),
   size=22, color=WHITE)
tb(sl, "Presentación Corporativa  ·  Abril 2026",
   Inches(1.10), Inches(4.05), Inches(7.0), Inches(0.40),
   size=12, color=MUTED)

# Badge SBS — tarjeta verde sobria
rect(sl, Inches(1.10), Inches(5.20), Inches(5.20), Inches(0.60), fill=GREEN_D)
tb(sl, "Registrada ante la SBS — Resolución N° 00313-2026",
   Inches(1.10), Inches(5.20), Inches(5.20), Inches(0.60),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "www.qoricash.pe",
   Inches(1.10), Inches(6.10), Inches(4.0), Inches(0.40),
   size=15, bold=True, color=GREEN)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — QUIÉNES SOMOS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Quiénes Somos",
                   "Identidad, misión y visión de QoriCash FX")

# ── Col izquierda: datos corporativos ──
CW = Inches(5.20)
card(sl, Inches(0.40), y0, CW, Inches(5.80), accent=True)
tb(sl, "Datos Corporativos",
   Inches(0.60), y0 + Inches(0.20), CW - Inches(0.30), Inches(0.40),
   size=13, bold=True, color=GREEN_D)

corp = [
    ("Razón Social",    "QORICASH S.A.C."),
    ("Nombre Comercial","QORICASH FX"),
    ("RUC",             "20615113698"),
    ("Registro SBS",    "Res. N° 00313-2026"),
    ("Dirección",       "Av. Brasil N° 2790, Int. 504 — Pueblo Libre, Lima"),
    ("Horario",         "Lun–Vie: 9:00am–6:00pm  ·  Sáb: 9:00am–1:00pm"),
    ("Contacto",        "info@qoricash.pe  ·  +51 926 011 920"),
]
ry = y0 + Inches(0.72)
for lbl, val in corp:
    tb(sl, lbl, Inches(0.60), ry, Inches(1.90), Inches(0.42),
       size=10, bold=True, color=NAVY)
    tb(sl, val, Inches(2.55), ry, Inches(2.85), Inches(0.42),
       size=10, color=BODY)
    rect(sl, Inches(0.60), ry + Inches(0.43),
         CW - Inches(0.35), Inches(0.01), fill=BORDER)
    ry += Inches(0.56)

# ── Col derecha: Misión / Visión / Valores ──
RX = Inches(6.00)
RW = Inches(6.90)

# Misión — fondo verde muy claro
rect(sl, RX, y0, RW, Inches(1.80), fill=GREEN_XL, line=GREEN_100, lw=Pt(0.5))
rect(sl, RX, y0, Inches(0.06), Inches(1.80), fill=GREEN_D)
icon_box(sl, "🎯", RX + Inches(0.18), y0 + Inches(0.18), Inches(0.50),
         bg=GREEN_D, fg=WHITE, font_size=14)
tb(sl, "Misión", RX + Inches(0.84), y0 + Inches(0.12), RW - Inches(1.0), Inches(0.40),
   size=13, bold=True, color=GREEN_D)
tb(sl, "Democratizar el acceso a tipos de cambio justos, eliminando las barreras "
       "del cambio de divisas mediante una plataforma digital segura, rápida y transparente.",
   RX + Inches(0.18), y0 + Inches(0.72), RW - Inches(0.35), Inches(0.96),
   size=11, color=BODY)

# Visión — fondo blanco con borde
my = y0 + Inches(1.95)
rect(sl, RX, my, RW, Inches(1.80), fill=WHITE, line=BORDER, lw=Pt(0.5))
rect(sl, RX, my, Inches(0.06), Inches(1.80), fill=GREEN_D)
icon_box(sl, "🏆", RX + Inches(0.18), my + Inches(0.18), Inches(0.50),
         bg=GREEN_100, fg=GREEN_D, font_size=14)
tb(sl, "Visión", RX + Inches(0.84), my + Inches(0.12), RW - Inches(1.0), Inches(0.40),
   size=13, bold=True, color=GREEN_D)
tb(sl, "Ser la casa de cambio digital líder en el Perú, reconocida por innovación "
       "tecnológica, excelencia en el servicio y compromiso inquebrantable con la "
       "seguridad de nuestros clientes.",
   RX + Inches(0.18), my + Inches(0.72), RW - Inches(0.35), Inches(0.96),
   size=11, color=BODY)

# Valores — 4 píldoras con icono
vy = my + Inches(1.95)
tb(sl, "Valores", RX + Inches(0.18), vy, RW - Inches(0.30), Inches(0.38),
   size=11, bold=True, color=NAVY)
vy += Inches(0.42)
valores = [("🛡", "Seguridad"), ("⚡", "Rapidez"), ("🤝", "Transparencia"), ("⭐", "Confianza")]
vx = RX
vw_pill = Inches(1.58)
for ico, vname in valores:
    rect(sl, vx, vy, vw_pill, Inches(0.68), fill=GREEN_100, line=GREEN, lw=Pt(0.5))
    icon_box(sl, ico, vx + Inches(0.08), vy + Inches(0.09), Inches(0.50),
             bg=GREEN_D, fg=WHITE, font_size=8)
    tb(sl, vname, vx + Inches(0.62), vy + Inches(0.18),
       vw_pill - Inches(0.70), Inches(0.36),
       size=10, bold=True, color=GREEN_D)
    vx += vw_pill + Inches(0.14)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — A QUÉ NOS DEDICAMOS (propuesta de valor)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "A Qué Nos Dedicamos",
                   "Cambio de divisas que rentabiliza el dinero de nuestros clientes")

# ── Bloque hero central: la esencia del negocio ──
rect(sl, Inches(0.40), y0, W - Inches(0.80), Inches(1.55), fill=NAVY)
rect(sl, Inches(0.40), y0, Inches(0.08), Inches(1.55), fill=GREEN)
# Icono FX decorativo en el lado derecho del hero
icon_box(sl, "💱", W - Inches(1.90), y0 + Inches(0.28), Inches(1.00),
         bg=GREEN_D, fg=WHITE, font_size=22)
tb(sl, "Compramos y vendemos dólares (USD) al mejor tipo de cambio del Perú,",
   Inches(0.65), y0 + Inches(0.18), Inches(10.80), Inches(0.50),
   size=17, bold=True, color=WHITE)
tb(sl, "maximizando el valor de cada sol o cada dólar que nuestros clientes necesitan convertir.",
   Inches(0.65), y0 + Inches(0.72), Inches(10.80), Inches(0.55),
   size=14, color=GREEN)

# ── Flujo visual USD ↔ PEN ──
FY = y0 + Inches(1.80)
FH = Inches(1.55)

rect(sl, Inches(0.40), FY, Inches(2.80), FH, fill=GREEN_XL, line=GREEN, lw=Pt(0.5))
rect(sl, Inches(0.40), FY, Inches(0.06), FH, fill=GREEN_D)
icon_box(sl, "💵", Inches(0.55), FY + Inches(0.14), Inches(0.60),
         bg=GREEN_D, fg=WHITE, font_size=14)
tb(sl, "CLIENTE", Inches(1.25), FY + Inches(0.14), Inches(1.85), Inches(0.36),
   size=10, bold=True, color=GREEN_D)
tb(sl, "Entrega\nUSD o PEN",
   Inches(0.55), FY + Inches(0.60), Inches(2.55), Inches(0.75),
   size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

tb(sl, "→", Inches(3.30), FY + Inches(0.50), Inches(0.60), Inches(0.55),
   size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

rect(sl, Inches(4.00), FY, Inches(4.30), FH, fill=NAVY)
rect(sl, Inches(4.00), FY, Inches(0.08), FH, fill=GREEN)
icon_box(sl, "🔄", Inches(4.18), FY + Inches(0.14), Inches(0.65),
         bg=GREEN_D, fg=WHITE, font_size=16)
tb(sl, "QoriCash FX", Inches(4.95), FY + Inches(0.14), Inches(3.25), Inches(0.36),
   size=13, bold=True, color=GREEN)
tb(sl, "Aplica el mejor\ntipo de cambio del mercado",
   Inches(4.18), FY + Inches(0.60), Inches(4.05), Inches(0.75),
   size=13, color=WHITE, align=PP_ALIGN.CENTER)

tb(sl, "→", Inches(8.40), FY + Inches(0.50), Inches(0.60), Inches(0.55),
   size=30, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

rect(sl, Inches(9.10), FY, Inches(3.80), FH, fill=GREEN_XL, line=GREEN, lw=Pt(0.5))
rect(sl, Inches(9.10), FY, Inches(0.06), FH, fill=GREEN_D)
icon_box(sl, "💰", Inches(9.25), FY + Inches(0.14), Inches(0.60),
         bg=GREEN_D, fg=WHITE, font_size=14)
tb(sl, "CLIENTE", Inches(9.95), FY + Inches(0.14), Inches(2.80), Inches(0.36),
   size=10, bold=True, color=GREEN_D)
tb(sl, "Recibe más\npor su dinero",
   Inches(9.25), FY + Inches(0.60), Inches(3.55), Inches(0.75),
   size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# ── 4 pilares de valor con iconos ──
PY = y0 + Inches(3.58)
pilares = [
    ("👥", "Para quién",
     "Personas naturales y empresas\n(RUC) con cuenta bancaria en Perú"),
    ("💱", "Qué cambiamos",
     "Dólares americanos (USD)\npor soles peruanos (PEN) y viceversa"),
    ("⚡", "Cómo lo hacemos",
     "100% digital, sin colas,\ntransferencia en menos de 15 minutos"),
    ("📈", "Por qué somos mejores",
     "Tasas sobre bancos, sin comisiones\nocultas, con factura electrónica SUNAT"),
]
PW = Inches(2.95); PH = Inches(2.55); PGAP = Inches(0.23)
for i, (ico, title, body) in enumerate(pilares):
    px = Inches(0.40) + i * (PW + PGAP)
    card(sl, px, PY, PW, PH, accent=True)
    icon_box(sl, ico, px + Inches(0.18), PY + Inches(0.18), Inches(0.56),
             bg=GREEN_D, fg=WHITE, font_size=9)
    tb(sl, title, px + Inches(0.84), PY + Inches(0.22),
       PW - Inches(0.98), Inches(0.42),
       size=12, bold=True, color=GREEN_D)
    tb(sl, body, px + Inches(0.18), PY + Inches(0.84),
       PW - Inches(0.32), Inches(1.55),
       size=11, color=BODY)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — CIFRAS QUE NOS RESPALDAN
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Cifras que nos Respaldan",
                   "Resultados reales, construidos operación a operación")

stats = [
    ("👥",  "8,500+",  "Usuarios\nregistrados"),
    ("💰",  "S/ 18M+", "Soles\ncambiados"),
    ("🔄",  "4,200+",  "Operaciones\ncompletadas"),
    ("⭐",  "4.8 / 5", "Satisfacción\nde clientes"),
]
sw = Inches(2.80); sh = Inches(4.30); gap = Inches(0.35)
sx0 = Inches(0.55)
for i, (ico, num, lbl) in enumerate(stats):
    sx = sx0 + i * (sw + gap)
    card(sl, sx, y0, sw, sh, accent=True)
    # Icono grande centrado en la parte superior
    icon_box(sl, ico, sx + (sw - Inches(1.10)) / 2, y0 + Inches(0.28),
             Inches(1.10), bg=GREEN_D, fg=WHITE, font_size=22)
    # Línea separadora
    rect(sl, sx + Inches(0.25), y0 + Inches(1.55), sw - Inches(0.50), Inches(0.01), fill=GREEN_100)
    tb(sl, num, sx, y0 + Inches(1.70), sw, Inches(1.00),
       size=40, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)
    tb(sl, lbl, sx, y0 + Inches(2.80), sw, Inches(0.90),
       size=13, color=MUTED, align=PP_ALIGN.CENTER)

# Tagline inferior
rect(sl, Inches(0.55), y0 + Inches(4.48), W - Inches(1.10), Inches(0.54), fill=GREEN_XL, line=GREEN_100, lw=Pt(0.3))
rect(sl, Inches(0.55), y0 + Inches(4.48), Inches(0.06), Inches(0.54), fill=GREEN)
tb(sl, "Más de 8,500 clientes ya eligieron el mejor tipo de cambio del Perú.",
   Inches(0.75), y0 + Inches(4.50), W - Inches(1.30), Inches(0.50),
   size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — CÓMO FUNCIONA
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Cómo Funciona",
                   "5 pasos simples para realizar su operación de cambio")

steps = [
    ("1", "📈",  "Cotización\nen tiempo real",
     "Consulte el tipo de cambio en\nwww.qoricash.pe o la app móvil."),
    ("2", "🪪",  "Registro\nKYC gratuito",
     "Cree su cuenta con DNI, CE\no RUC. 100% digital."),
    ("3", "📝",  "Solicite\nsu operación",
     "Indique monto y cuentas\nbancarias origen y destino."),
    ("4", "💸",  "Realice la\ntransferencia",
     "Transfiera a QoriCash\ny adjunte su comprobante."),
    ("5", "✅",  "Reciba\nsu dinero",
     "En menos de 15 min recibe\nel monto y la factura SUNAT."),
]
sw2 = Inches(2.28); sh2 = Inches(4.85); gap2 = Inches(0.12)
for i, (num, ico, title, desc) in enumerate(steps):
    sx = Inches(0.40) + i * (sw2 + gap2)
    card(sl, sx, y0, sw2, sh2, accent=False)
    # Barra top verde
    rect(sl, sx, y0, sw2, Inches(0.07), fill=GREEN)
    # Número en círculo verde
    rect(sl, sx + (sw2 - Inches(0.80)) / 2, y0 + Inches(0.20),
         Inches(0.80), Inches(0.80), fill=GREEN)
    tb(sl, num, sx + (sw2 - Inches(0.80)) / 2, y0 + Inches(0.20),
       Inches(0.80), Inches(0.80),
       size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Icono del paso
    icon_box(sl, ico, sx + (sw2 - Inches(0.65)) / 2, y0 + Inches(1.18),
             Inches(0.65), bg=GREEN_XL, fg=GREEN_D, font_size=11)
    # Título
    tb(sl, title, sx + Inches(0.10), y0 + Inches(2.04),
       sw2 - Inches(0.20), Inches(0.85),
       size=13, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    # Separador
    rect(sl, sx + Inches(0.30), y0 + Inches(2.94),
         sw2 - Inches(0.60), Inches(0.01), fill=BORDER)
    # Descripción
    tb(sl, desc, sx + Inches(0.12), y0 + Inches(3.05),
       sw2 - Inches(0.24), Inches(1.60),
       size=10, color=MUTED, align=PP_ALIGN.CENTER)
    # Flecha entre pasos
    if i < len(steps) - 1:
        ax = sx + sw2 + Inches(0.01)
        tb(sl, "›", ax, y0 + Inches(1.90),
           Inches(0.12), Inches(0.60),
           size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — CANALES Y BANCOS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Canales y Cobertura Bancaria",
                   "Múltiples formas de acceder a QoriCash")

# ── Col izquierda ──
CL = Inches(5.60)
tb(sl, "Canales de Acceso",
   Inches(0.40), y0 + Inches(0.05), CL, Inches(0.38),
   size=13, bold=True, color=GREEN_D)

channels = [
    ("🌐", "Plataforma Web",  "www.qoricash.pe — disponible desde cualquier navegador"),
    ("📱", "App Móvil",       "iOS y Android — opere desde su celular"),
    ("📊", "Dashboard",       "Historial, cuentas y seguimiento en tiempo real"),
    ("💬", "WhatsApp",        "+51 926 011 920 — soporte por operadores humanos"),
]
cy2 = y0 + Inches(0.52)
for icon_ch, ch, desc in channels:
    card(sl, Inches(0.40), cy2, CL, Inches(1.02), accent=False)
    rect(sl, Inches(0.40), cy2, Inches(0.06), Inches(1.02), fill=GREEN)
    icon_box(sl, icon_ch, Inches(0.56), cy2 + Inches(0.24),
             Inches(0.54), bg=GREEN_D, fg=WHITE, font_size=11)
    tb(sl, ch, Inches(1.22), cy2 + Inches(0.10),
       CL - Inches(0.92), Inches(0.38),
       size=12, bold=True, color=NAVY)
    tb(sl, desc, Inches(1.22), cy2 + Inches(0.50),
       CL - Inches(0.92), Inches(0.42),
       size=10, color=MUTED)
    cy2 += Inches(1.14)

# ── Col derecha — Bancos ──
BX = Inches(6.40); BW = Inches(6.50)
tb(sl, "Entidades Bancarias",
   BX, y0 + Inches(0.05), BW, Inches(0.38),
   size=13, bold=True, color=GREEN_D)

# Inmediatas
rect(sl, BX, y0 + Inches(0.52), BW, Inches(0.44), fill=GREEN_D)
tb(sl, "TRANSFERENCIAS INMEDIATAS — Todo el Perú",
   BX, y0 + Inches(0.52), BW, Inches(0.44),
   size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

iy = y0 + Inches(1.06)
for b in ["BCP", "Interbank"]:
    rect(sl, BX, iy, BW, Inches(0.52), fill=OFF_WHITE, line=BORDER, lw=Pt(0.3))
    icon_box(sl, b[:2], BX + Inches(0.12), iy + Inches(0.09),
             Inches(0.34), bg=GREEN_100, fg=GREEN_D, font_size=7)
    tb(sl, b, BX + Inches(0.55), iy + Inches(0.10),
       Inches(3.70), Inches(0.34), size=12, color=BODY)
    rect(sl, BX + Inches(4.70), iy + Inches(0.10),
         Inches(1.50), Inches(0.32), fill=GREEN_100)
    tb(sl, "Inmediato", BX + Inches(4.70), iy + Inches(0.10),
       Inches(1.50), Inches(0.32),
       size=10, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)
    iy += Inches(0.58)

# Interbancarias
rect(sl, BX, iy + Inches(0.18), BW, Inches(0.44),
     fill=CARD_BG, line=BORDER, lw=Pt(0.3))
tb(sl, "INTERBANCARIAS — Lima",
   BX, iy + Inches(0.18), BW, Inches(0.44),
   size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

iy2 = iy + Inches(0.74)
for b in ["BBVA", "Scotiabank", "BanBif", "Banco Pichincha", "Otros bancos"]:
    rect(sl, BX, iy2, BW, Inches(0.48), fill=WHITE, line=BORDER, lw=Pt(0.3))
    icon_box(sl, b[:2], BX + Inches(0.12), iy2 + Inches(0.09),
             Inches(0.30), bg=CARD_BG, fg=MUTED, font_size=6)
    tb(sl, b, BX + Inches(0.52), iy2 + Inches(0.09),
       Inches(3.70), Inches(0.32), size=11, color=BODY)
    tb(sl, "~ 20 min", BX + Inches(4.80), iy2 + Inches(0.09),
       Inches(1.40), Inches(0.32),
       size=10, color=MUTED, align=PP_ALIGN.CENTER)
    iy2 += Inches(0.52)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — BENEFICIOS PARA EL CLIENTE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Beneficios para el Cliente",
                   "Por qué cambiar con QoriCash es siempre mejor")

benefits = [
    ("📈",  "Mejor Tipo de Cambio",
     "Tasas actualizadas en tiempo real,\nsiempre superiores a los bancos del sistema."),
    ("🤝",  "Sin Comisiones Ocultas",
     "El tipo de cambio que ve es exactamente\nel que obtiene. Sin costos adicionales."),
    ("⚡",  "Rapidez Garantizada",
     "BCP / Interbank: inmediato.\nResto del sistema bancario: < 15 minutos."),
    ("📱",  "100% Digital",
     "Sin colas. Opere desde su celular\no computadora, las 24 horas."),
    ("📄",  "Factura Electrónica",
     "Boleta o factura SUNAT automática\npara cada operación completada."),
    ("🔗",  "Programa de Referidos",
     "Acumule pips por referidos y mejore\nsu tasa efectiva en futuras operaciones."),
]
bw3 = Inches(3.90); bh3 = Inches(2.45); gap3 = Inches(0.20)
for i, (icon_b, title, desc) in enumerate(benefits):
    col = i % 3; row = i // 3
    bx = Inches(0.40) + col * (bw3 + gap3)
    by = y0 + Inches(0.05) + row * (bh3 + Inches(0.10))
    card(sl, bx, by, bw3, bh3, accent=True)
    # Icono verde
    icon_box(sl, icon_b, bx + Inches(0.20), by + Inches(0.22),
             Inches(0.62), bg=GREEN_D, fg=WHITE, font_size=10)
    # Título + separador
    tb(sl, title, bx + Inches(0.96), by + Inches(0.22),
       bw3 - Inches(1.10), Inches(0.46),
       size=12, bold=True, color=NAVY)
    rect(sl, bx + Inches(0.20), by + Inches(0.82),
         bw3 - Inches(0.40), Inches(0.01), fill=BORDER)
    tb(sl, desc, bx + Inches(0.20), by + Inches(0.94),
       bw3 - Inches(0.40), Inches(1.38),
       size=11, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — TIPO DE CAMBIO COMPARATIVO
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "¿Cuánto Ahorra con QoriCash?",
                   "Comparamos nuestras tasas vs. los principales bancos del Perú")

# Panel izquierdo highlight
HW = Inches(4.40)
rect(sl, Inches(0.40), y0, HW, Inches(5.50), fill=GREEN_XL, line=GREEN_100, lw=Pt(0.5))
rect(sl, Inches(0.40), y0, Inches(0.06), Inches(5.50), fill=GREEN)

# Icono central grande
icon_box(sl, "💰", Inches(0.40) + (HW - Inches(1.40)) / 2, y0 + Inches(0.28),
         Inches(1.40), bg=GREEN_D, fg=WHITE, font_size=28)

tb(sl, "Al cambiar $1,000 USD",
   Inches(0.55), y0 + Inches(1.88), HW - Inches(0.25), Inches(0.40),
   size=11, color=GREEN_D, align=PP_ALIGN.CENTER)
tb(sl, "Hasta S/ 100",
   Inches(0.40), y0 + Inches(2.32), HW, Inches(1.00),
   size=44, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
tb(sl, "más que en un banco",
   Inches(0.40), y0 + Inches(3.32), HW, Inches(0.40),
   size=14, color=BODY, align=PP_ALIGN.CENTER)
rect(sl, Inches(0.90), y0 + Inches(3.84), Inches(3.40), Inches(0.01), fill=GREEN_100)
tb(sl, "La diferencia crece con el volumen.\nA mayor monto, mayor ahorro.",
   Inches(0.60), y0 + Inches(3.96), HW - Inches(0.30), Inches(0.78),
   size=10, color=MUTED, align=PP_ALIGN.CENTER)

# Tabla comparativa derecha
TX = Inches(5.20); TW = [Inches(3.10), Inches(1.65), Inches(1.65)]
headers = ["Entidad", "Compra (S/)", "Venta (S/)"]
rows = [
    ("QoriCash FX  ★", "MEJOR TASA", "MEJOR TASA", True),
    ("BCP",             "3.340",       "3.440",       False),
    ("Interbank",       "3.345",       "3.435",       False),
    ("BBVA",            "3.330",       "3.450",       False),
    ("Scotiabank",      "3.325",       "3.455",       False),
]
ty2 = y0
rx = TX
for hdr, tw in zip(headers, TW):
    rect(sl, rx, ty2, tw, Inches(0.50), fill=NAVY)
    tb(sl, hdr, rx, ty2, tw, Inches(0.50),
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rx += tw
ty2 += Inches(0.50)

for ent, comp, vent, highlight in rows:
    bg = GREEN_XL if highlight else (OFF_WHITE if rows.index((ent,comp,vent,highlight)) % 2 == 1 else WHITE)
    rx = TX
    for j, (val, tw) in enumerate(zip([ent, comp, vent], TW)):
        rect(sl, rx, ty2, tw, Inches(0.54),
             fill=bg, line=BORDER, lw=Pt(0.3))
        fc = GREEN_D if highlight else BODY
        tb(sl, val, rx + Inches(0.12), ty2 + Inches(0.10),
           tw - Inches(0.12), Inches(0.36),
           size=11, bold=highlight, color=fc,
           align=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)
        rx += tw
    ty2 += Inches(0.54)

tb(sl, "* Tasas bancarias son referenciales. Las de QoriCash se actualizan en tiempo real durante el horario de atención.",
   TX, ty2 + Inches(0.12), Inches(6.40), Inches(0.40),
   size=9, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — SEGURIDAD Y COMPLIANCE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Seguridad y Cumplimiento Regulatorio",
                   "La confianza no se declara — se construye con hechos")

pillars = [
    ("⚖",  "Registro SBS",          "Res. N° 00313-2026",
     "Operamos bajo supervisión de la SBS. Empresa legalmente "
     "constituida con plena responsabilidad regulatoria ante el Estado peruano."),
    ("🪪",  "KYC",                   "Verificación de identidad",
     "Todo cliente verifica su DNI, CE o RUC antes de operar. "
     "Documentos almacenados en entornos seguros en la nube."),
    ("🔍",  "PLAFT / AML",           "Prevención de Lavado de Activos",
     "Score de riesgo 0–100, monitoreo continuo, alertas por monto "
     "y frecuencia, screening contra listas internacionales de sanciones y PEPs."),
    ("🔒",  "Cifrado y Privacidad",  "Protección de datos",
     "Cifrado de extremo a extremo. Acceso restringido y trazabilidad "
     "completa de cada acción realizada en el sistema."),
]
pw2 = Inches(2.90); ph2 = Inches(5.05); gap4 = Inches(0.25)
for i, (icon_p, title, badge, desc) in enumerate(pillars):
    px = Inches(0.40) + i * (pw2 + gap4)
    card(sl, px, y0, pw2, ph2, accent=True)
    # Icono grande centrado en parte superior
    icon_box(sl, icon_p, px + (pw2 - Inches(1.00)) / 2, y0 + Inches(0.22),
             Inches(1.00), bg=GREEN_D, fg=WHITE, font_size=16)
    # Badge verde pálido
    rect(sl, px + Inches(0.20), y0 + Inches(1.40),
         pw2 - Inches(0.40), Inches(0.36), fill=GREEN_100, line=GREEN, lw=Pt(0.3))
    tb(sl, badge, px + Inches(0.20), y0 + Inches(1.40),
       pw2 - Inches(0.40), Inches(0.36),
       size=9, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)
    # Separador
    rect(sl, px + Inches(0.20), y0 + Inches(1.88),
         pw2 - Inches(0.40), Inches(0.01), fill=BORDER)
    # Título
    tb(sl, title, px + Inches(0.18), y0 + Inches(2.00),
       pw2 - Inches(0.36), Inches(0.50),
       size=14, bold=True, color=NAVY)
    # Descripción
    tb(sl, desc, px + Inches(0.18), y0 + Inches(2.58),
       pw2 - Inches(0.36), Inches(2.30),
       size=10, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — 8 RAZONES
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "8 Razones para Elegir QoriCash FX",
                   "Atributos que nos posicionan como líderes del mercado")

reasons = [
    ("01", "⚖",  "Supervisados por la SBS",    "Registro oficial Res. 00313-2026."),
    ("02", "📈",  "Mejor tipo de cambio",        "Siempre por encima de los bancos."),
    ("03", "💻",  "Tecnología de punta",         "Web + App móvil + operadores humanos."),
    ("04", "💬",  "Atención multicanal real",    "WhatsApp + plataforma + dashboard."),
    ("05", "🏦",  "Cobertura bancaria amplia",   "7 entidades del sistema financiero peruano."),
    ("06", "🏢",  "Gestión empresarial",         "RUC, facturas electrónicas y múltiples cuentas."),
    ("07", "✅",  "Trazabilidad total",          "Historial, comprobantes y seguimiento en vivo."),
    ("08", "🔗",  "Programa de fidelización",    "Referidos con beneficios directos en la tasa."),
]
rw2 = Inches(2.90); rh2 = Inches(1.50); gap5 = Inches(0.20)
cols3 = 4
for i, (num, ico, title, desc) in enumerate(reasons):
    col = i % cols3; row = i // cols3
    rx2 = Inches(0.40) + col * (rw2 + gap5)
    ry2 = y0 + Inches(0.05) + row * (rh2 + Inches(0.16))
    card(sl, rx2, ry2, rw2, rh2, accent=False)
    rect(sl, rx2, ry2, Inches(0.06), rh2, fill=GREEN)
    # Número + icono del concepto
    icon_box(sl, num, rx2 + Inches(0.14), ry2 + Inches(0.12),
             Inches(0.50), bg=GREEN_XL, fg=GREEN_D, font_size=9)
    icon_box(sl, ico, rx2 + rw2 - Inches(0.64), ry2 + Inches(0.12),
             Inches(0.50), bg=CARD_BG, fg=NAVY, font_size=8)
    # Título
    tb(sl, title, rx2 + Inches(0.76), ry2 + Inches(0.12),
       rw2 - Inches(1.40), Inches(0.46),
       size=11, bold=True, color=NAVY)
    # Separador
    rect(sl, rx2 + Inches(0.14), ry2 + Inches(0.68),
         rw2 - Inches(0.22), Inches(0.01), fill=BORDER)
    # Desc
    tb(sl, desc, rx2 + Inches(0.14), ry2 + Inches(0.76),
       rw2 - Inches(0.24), Inches(0.58),
       size=9, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — SOCIOS ESTRATÉGICOS
# ═══════════════════════════════════════════════════════════════════════════
LOGOS_DIR = "/Users/gianpierre/Desktop/Qoricash/Sistema/qoricash/docs/logos/"

sl = prs.slides.add_slide(BLANK)
y0 = section_title(sl, "Socios y Alianzas Estratégicas",
                   "Trabajamos de la mano con las instituciones que impulsan el ecosistema financiero del Perú")

socios = [
    ("ADEX",  LOGOS_DIR + "adex.png",
     "Asociación de Exportadores del Perú",
     "Apoyamos a empresas exportadoras e importadoras en la gestión eficiente de sus operaciones de cambio de divisas."),
    ("CCL",   LOGOS_DIR + "ccl.jpg",
     "Cámara de Comercio de Lima",
     "Alianza con el principal referente empresarial del Perú para llevar soluciones de cambio competitivo al sector comercial."),
    ("BCRP",  LOGOS_DIR + "bcrp.png",
     "Banco Central de Reserva del Perú",
     "Operamos alineados al marco regulatorio y las publicaciones de tipo de cambio de referencia del BCRP."),
    ("BVL",   LOGOS_DIR + "bvl.png",
     "Bolsa de Valores de Lima",
     "Alianza estratégica con el mercado de valores peruano para ofrecer soluciones de cambio al ecosistema bursátil e inversores."),
]

SW2 = Inches(2.95)
SH2 = Inches(4.90)
SY2 = y0 + Inches(0.10)
SGAP2 = Inches(0.26)

for i, (name, logo_path, org, desc) in enumerate(socios):
    sx = Inches(0.40) + i * (SW2 + SGAP2)
    card(sl, sx, SY2, SW2, SH2, accent=False)
    rect(sl, sx, SY2, SW2, Inches(0.06), fill=GREEN)

    # Área de logo
    LOGO_ZONE_H = Inches(1.50)
    rect(sl, sx + Inches(0.12), SY2 + Inches(0.18),
         SW2 - Inches(0.24), LOGO_ZONE_H, fill=OFF_WHITE)

    if logo_path:
        try:
            pic_h = Inches(0.96)
            pic_y = SY2 + Inches(0.18) + (LOGO_ZONE_H - pic_h) / 2
            sl.shapes.add_picture(logo_path,
                                  sx + Inches(0.20), pic_y,
                                  width=SW2 - Inches(0.40), height=pic_h)
        except Exception:
            tb(sl, name, sx + Inches(0.12), SY2 + Inches(0.60),
               SW2 - Inches(0.24), Inches(0.55),
               size=18, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)
    else:
        tb(sl, name, sx + Inches(0.12), SY2 + Inches(0.60),
           SW2 - Inches(0.24), Inches(0.55),
           size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Separador
    rect(sl, sx + Inches(0.20), SY2 + LOGO_ZONE_H + Inches(0.22),
         SW2 - Inches(0.40), Inches(0.01), fill=GREEN)

    # Nombre organización
    tb(sl, org,
       sx + Inches(0.14), SY2 + LOGO_ZONE_H + Inches(0.30),
       SW2 - Inches(0.28), Inches(0.64),
       size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

    # Separador 2
    rect(sl, sx + Inches(0.30), SY2 + LOGO_ZONE_H + Inches(1.00),
         SW2 - Inches(0.60), Inches(0.01), fill=GREEN_100)

    # Descripción
    tb(sl, desc,
       sx + Inches(0.14), SY2 + LOGO_ZONE_H + Inches(1.10),
       SW2 - Inches(0.28), Inches(2.50),
       size=9, color=MUTED, align=PP_ALIGN.CENTER)

# Tagline inferior
rect(sl, Inches(0.40), y0 + Inches(5.22), W - Inches(0.80), Inches(0.52), fill=GREEN_XL)
rect(sl, Inches(0.40), y0 + Inches(5.22), Inches(0.06), Inches(0.52), fill=GREEN)
tb(sl, "Juntos construimos un ecosistema financiero más transparente, competitivo y accesible para el Perú.",
   Inches(0.60), y0 + Inches(5.24), W - Inches(1.20), Inches(0.48),
   size=11, bold=True, color=GREEN_D, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — CIERRE Y CONTACTO  (navy — simétrico con portada)
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, W, H, fill=NAVY)
rect(sl, 0, 0, Inches(0.50), H, fill=GREEN)
rect(sl, Inches(0.50), H - Inches(0.10), W - Inches(0.50), Inches(0.10), fill=GREEN_D)
# Decoración geométrica: bloques en la esquina derecha
for k, (bx4, bw4, ba) in enumerate([
    (W - Inches(0.80), Inches(0.80), 20),
    (W - Inches(2.10), Inches(1.10), 10),
    (W - Inches(3.80), Inches(1.40), 5),
]):
    blk = rect(sl, bx4, 0, bw4, H, fill=GREEN_D)
    set_shape_fill_alpha(blk, ba)
# Líneas forex sutiles de fondo
for k in range(1, 6):
    fline(sl, Inches(0.50), H * k / 6, W, H * k / 6, GREEN, width_pt=0.3, alpha=4)

logo(sl, Inches(5.40), Inches(0.40), Inches(1.80))

tb(sl, "QoriCash FX",
   Inches(0.80), Inches(2.10), W - Inches(1.60), Inches(0.90),
   size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(sl, "Créanos hoy. Compruébalo en su primera operación.",
   Inches(0.80), Inches(3.05), W - Inches(1.60), Inches(0.52),
   size=17, color=GREEN, align=PP_ALIGN.CENTER)

rect(sl, Inches(2.0), Inches(3.72), Inches(9.30), Inches(0.015), fill=GREEN_D)

contacts = [
    ("🌐",  "Sitio Web",  "www.qoricash.pe"),
    ("📞",  "WhatsApp",   "+51 926 011 920"),
    ("✉",   "Email",      "info@qoricash.pe"),
    ("📍",  "Dirección",  "Av. Brasil N° 2790, Int. 504 — Pueblo Libre, Lima"),
    ("🕐",  "Horario",    "Lun–Vie: 9am–6pm  ·  Sáb: 9am–1pm"),
]
cy3 = Inches(3.92)
for ico3, lbl, val in contacts:
    icon_box(sl, ico3, Inches(2.30), cy3, Inches(0.38),
             bg=GREEN_D, fg=WHITE, font_size=6)
    tb(sl, lbl + ":", Inches(2.78), cy3, Inches(1.80), Inches(0.40),
       size=11, bold=True, color=GREEN, align=PP_ALIGN.LEFT)
    tb(sl, val, Inches(4.68), cy3, Inches(7.20), Inches(0.40),
       size=11, color=WHITE)
    cy3 += Inches(0.46)

rect(sl, Inches(0.60), Inches(6.55), Inches(12.10), Inches(0.48),
     fill=RGBColor(0x05, 0x10, 0x1A))
tb(sl, "QORICASH S.A.C.  ·  RUC 20615113698  ·  SBS Res. N° 00313-2026  ·  © 2026 QoriCash FX",
   Inches(0.60), Inches(6.55), Inches(12.10), Inches(0.48),
   size=8, color=MUTED, align=PP_ALIGN.CENTER)


# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"OK — Presentación guardada en:\n{OUT}")
